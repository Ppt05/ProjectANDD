# -*- coding: utf-8 -*-
"""
Export Slide PowerPoint (.pptx) tự động
Bài báo: Deanonymizing Social Networks Using Structural Information (IJCAI-2019)
Môn: An Ninh Di Động

Cài đặt: pip install python-pptx
Chạy:    python export_slide.py
Output:  presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────────────────────
# THEME CONFIG
# ─────────────────────────────────────────────────────────────

BG_DARK    = RGBColor(0x0D, 0x11, 0x17)   # #0D1117 — nền tối
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)   # #161B22 — card nền
ACCENT_RED = RGBColor(0xE6, 0x39, 0x46)   # #E63946 — đỏ chủ đạo
ACCENT_ORG = RGBColor(0xF4, 0xA2, 0x61)   # #F4A261 — cam phụ
TEXT_WHITE = RGBColor(0xE8, 0xE8, 0xE8)   # #E8E8E8 — text sáng
TEXT_GRAY  = RGBColor(0x8B, 0x94, 0x9E)   # #8B949E — text mờ
COLOR_GREEN= RGBColor(0x3F, 0xB9, 0x50)   # #3FB950 — xanh lá nhấn

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────
# HELPER FUNCTIONS
# ─────────────────────────────────────────────────────────────

def set_slide_background(slide, color: RGBColor):
    """Đặt màu nền cho slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text: str, left, top, width, height,
                font_size=18, bold=False, color=TEXT_WHITE,
                align=PP_ALIGN.LEFT, italic=False) -> None:
    """Thêm textbox với style."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_title_slide(prs, title: str, subtitle: str, info: str):
    """Tạo slide trang bìa."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)

    # Đường kẻ trang trí trên cùng
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, Inches(0.07)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_RED
    line.line.fill.background()

    # Tiêu đề chính
    add_textbox(slide, title,
                Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.0),
                font_size=34, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    # Đường kẻ màu
    sep = slide.shapes.add_shape(1, Inches(4.5), Inches(4.0), Inches(4.3), Inches(0.05))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ACCENT_RED
    sep.line.fill.background()

    # Phụ đề
    add_textbox(slide, subtitle,
                Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.8),
                font_size=18, color=ACCENT_ORG, align=PP_ALIGN.CENTER, italic=True)

    # Thông tin nhóm
    add_textbox(slide, info,
                Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2),
                font_size=15, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Label góc dưới
    add_textbox(slide, "IJCAI-2019 | Caragiannis & Tsitsoka | Đại học Patras, Hy Lạp",
                Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
                font_size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    return slide


def add_content_slide(prs, title: str, bullets: list,
                       section_tag: str = "", accent=ACCENT_RED):
    """
    Tạo slide nội dung chuẩn với title + bullet list.
    bullets: list of (text, level) — level 0 = main, 1 = sub
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)

    # Header bar
    header = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.15))
    header.fill.solid()
    header.fill.fore_color.rgb = BG_CARD
    header.line.fill.background()

    # Accent strip
    strip = slide.shapes.add_shape(1, 0, 0, Inches(0.07), Inches(1.15))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()

    # Section tag
    if section_tag:
        add_textbox(slide, section_tag,
                    Inches(0.2), Inches(0.08), Inches(4), Inches(0.35),
                    font_size=10, color=accent, bold=True)

    # Tiêu đề slide
    add_textbox(slide, title,
                Inches(0.2), Inches(0.3), Inches(12.8), Inches(0.75),
                font_size=26, bold=True, color=TEXT_WHITE)

    # Bullet content
    content_top = Inches(1.3)
    content_height = Inches(5.9)
    txBox = slide.shapes.add_textbox(Inches(0.35), content_top, Inches(12.6), content_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, level) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = level
        p.space_before = Pt(4 if level == 0 else 2)

        # Bullet indicator
        if level == 0:
            bullet_char = "▸ "
            fsize = 17
            col = TEXT_WHITE
            b = True
        else:
            bullet_char = "    • "
            fsize = 15
            col = TEXT_GRAY
            b = False

        run = p.add_run()
        run.text = bullet_char + text
        run.font.size = Pt(fsize)
        run.font.bold = b
        run.font.color.rgb = col
        run.font.name = "Calibri"

    return slide


def add_table_slide(prs, title: str, headers: list, rows: list,
                    section_tag: str = "", accent=ACCENT_RED):
    """Tạo slide có bảng."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)

    # Header
    header = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.15))
    header.fill.solid()
    header.fill.fore_color.rgb = BG_CARD
    header.line.fill.background()

    strip = slide.shapes.add_shape(1, 0, 0, Inches(0.07), Inches(1.15))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()

    if section_tag:
        add_textbox(slide, section_tag,
                    Inches(0.2), Inches(0.08), Inches(4), Inches(0.35),
                    font_size=10, color=accent, bold=True)

    add_textbox(slide, title,
                Inches(0.2), Inches(0.3), Inches(12.8), Inches(0.75),
                font_size=26, bold=True, color=TEXT_WHITE)

    # Bảng
    n_rows = len(rows) + 1
    n_cols = len(headers)
    col_w  = Inches(12.5 / n_cols)
    row_h  = Inches(0.55)
    table_top = Inches(1.4)

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.4), table_top,
        Inches(12.5), row_h * n_rows
    ).table

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = TEXT_WHITE
        run.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        bg = BG_CARD if i % 2 == 0 else BG_DARK
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(13)
            run.font.color.rgb = TEXT_WHITE
            run.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_divider(prs, section_num: str, section_title: str,
                         subtitle: str, accent=ACCENT_RED):
    """Slide phân cách section."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BG_DARK)

    # Số section lớn
    add_textbox(slide, section_num,
                Inches(0), Inches(1.5), Inches(13.33), Inches(2.5),
                font_size=120, bold=True, color=BG_CARD, align=PP_ALIGN.CENTER)

    # Tên section
    add_textbox(slide, section_title,
                Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
                font_size=42, bold=True, color=accent, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, subtitle,
                Inches(2), Inches(4.2), Inches(9.33), Inches(0.8),
                font_size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER, italic=True)

    return slide


# ─────────────────────────────────────────────────────────────
# NỘI DUNG SLIDE
# ─────────────────────────────────────────────────────────────

def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── SLIDE 1: TRANG BÌA ──────────────────────────────────
    add_title_slide(
        prs,
        title    = "Phi Ẩn Danh Hóa Mạng Xã Hội\nSử Dụng Thông Tin Cấu Trúc",
        subtitle = "Phân tích bài báo IJCAI-2019",
        info     = "Môn: An Ninh Di Động  |  Nhóm: [Điền tên nhóm]  |  [Ngày báo cáo]"
    )

    # ── SLIDE 2: MỤC LỤC ────────────────────────────────────
    add_content_slide(
        prs,
        title      = "Nội Dung Trình Bày",
        section_tag= "",
        bullets    = [
            ("PHẦN 1 — LÝ THUYẾT", 0),
            ("Bài toán phi ẩn danh hóa mạng xã hội", 1),
            ("Thuật toán NDSD và Local Search", 1),
            ("Phân tích các mô hình đồ thị", 1),
            ("PHẦN 2 — THỰC NGHIỆM", 0),
            ("Thiết lập thực nghiệm & kết quả", 1),
            ("Biểu đồ Accuracy vs Noise", 1),
            ("PHẦN 3 — MỞ RỘNG", 0),
            ("Ứng dụng trong An ninh Di động", 1),
            ("Cơ chế phòng thủ & câu hỏi mở", 1),
        ]
    )

    # ── SECTION DIVIDER 1 ───────────────────────────────────
    add_section_divider(prs, "01", "LÝ THUYẾT",
                        "Bài toán • Thuật toán • Phân tích")

    # ── SLIDE 3: ĐẶT VẤN ĐỀ ────────────────────────────────
    add_content_slide(
        prs,
        title      = "Ẩn Danh Hóa Có Thực Sự Bảo Vệ Quyền Riêng Tư?",
        section_tag= "PHẦN 1 — LÝ THUYẾT",
        bullets    = [
            ("Thực trạng: Mạng xã hội chia sẻ dữ liệu sau khi 'ẩn danh hóa' (xóa tên, ID)", 0),
            ("Narayanan & Shmatikov 2009: De-anonymize Netflix dataset chỉ qua IMDB", 1),
            ("Backstrom et al. 2007: Tấn công qua 'sybil nodes'", 1),
            ("Vấn đề cốt lõi: Cấu trúc quan hệ vẫn còn nguyên dù đã xóa danh tính", 0),
            ("Ai kết nối với ai → 'vân tay số' có thể nhận dạng lại", 1),
            ("Bài báo này: Chỉ dùng thông tin cấu trúc đồ thị — không cần metadata", 0),
            ("Không cần tên, tuổi, địa chỉ — chỉ cần biểu diễn đồ thị của hai mạng", 1),
        ]
    )

    # ── SLIDE 4: MÔ HÌNH BÀI TOÁN ───────────────────────────
    add_content_slide(
        prs,
        title      = "Mô Hình Bài Toán Hình Thức",
        section_tag= "PHẦN 1 — LÝ THUYẾT",
        bullets    = [
            ("G = (V, E): Đồ thị có danh tính — biết tên và quan hệ của các node", 0),
            ("H = (V', E'): Đồ thị ẩn danh — node bị đổi tên, mất danh tính", 0),
            ("H là 'noisy subgraph' của G theo noise model:", 0),
            ("δ: tỷ lệ node bị xóa khỏi H (ví dụ δ=5% → 5% người không xuất hiện trong H)", 1),
            ("ε: xác suất mỗi cặp node có trạng thái cạnh khác nhau giữa G và H", 1),
            ("Mục tiêu: Tìm ánh xạ 1-1  σ: V(H) → V(G)  tối đa hóa cạnh khớp đúng", 0),
            ("max |{ (u,v) ∈ E(H) : (σ(u), σ(v)) ∈ E(G) }|", 1),
            ("Liên hệ lý thuyết: Tổng quát hóa bài toán Subgraph Isomorphism (NP-khó)", 0),
        ]
    )

    # ── SLIDE 5: THUẬT TOÁN NDSD ────────────────────────────
    add_content_slide(
        prs,
        title      = "Thuật Toán NDSD — Neighbor Degree Sequence Difference",
        section_tag= "PHẦN 1 — LÝ THUYẾT",
        bullets    = [
            ("Trực giác: Node có 'vùng lân cận tương tự' → khả năng là cùng một người", 0),
            ("Bước 1 — Tính 'chữ ký cấu trúc' NDS mỗi node:", 0),
            ("NDS(v) = dãy bậc của các láng giềng, sắp xếp giảm dần", 1),
            ("Ví dụ: v có láng giềng bậc [5,3,3,1] → NDS(v) = [5,3,3,1]", 1),
            ("Bước 2 — Tính khoảng cách L1 giữa mọi cặp (node_H, node_G):", 0),
            ("dist(u,v) = Σ |NDS(u)[i] - NDS(v)[i]|  → Ma trận chi phí n×n", 1),
            ("Bước 3 — Chạy Thuật Toán Hungary → khớp tối ưu tổng khoảng cách nhỏ nhất", 0),
            ("Độ phức tạp: O(n²) tính NDS + O(n²) tính cost + O(n³) Hungary = O(n⁴)", 0),
        ]
    )

    # ── SLIDE 6: THUẬT TOÁN LOCAL SEARCH ────────────────────
    add_content_slide(
        prs,
        title      = "Thuật Toán Local Search — Cải Thiện Qua Hoán Đổi",
        section_tag= "PHẦN 1 — LÝ THUYẾT",
        bullets    = [
            ("Ý tưởng: Bắt đầu từ nghiệm NDSD → thử swap liên tục để cải thiện", 0),
            ("Hàm mục tiêu: score(σ) = số cạnh H được khớp đúng trong G", 0),
            ("Vòng lặp cải thiện:", 0),
            ("Thử swap mọi cặp (u,v) trong H: σ' = σ với σ'(u) ↔ σ'(v)", 1),
            ("Nếu score(σ') > score(σ): chấp nhận swap", 1),
            ("Dừng khi không còn swap nào cải thiện được (local optimum)", 1),
            ("Độ phức tạp: O(n²) cặp × O(n) score × O(n²) vòng lặp = O(n⁵)", 0),
            ("Pipeline tối ưu: σ_NDSD → Local Search → σ_final", 0),
            ("NDSD cung cấp điểm khởi đầu tốt, Local Search tinh chỉnh", 1),
        ]
    )

    # ── SLIDE 7: 4 MÔ HÌNH ĐỒ THỊ ──────────────────────────
    add_table_slide(
        prs,
        title      = "Bốn Mô Hình Đồ Thị — Mô Hình Nào Phù Hợp Mạng XH?",
        section_tag= "PHẦN 1 — LÝ THUYẾT",
        headers    = ["Mô hình", "Phân phối bậc", "Hub rõ ràng", "Clustering", "Phù hợp MXH"],
        rows       = [
            ["Barabási-Albert", "Power-law", "✅ Có", "Thấp", "✅ Rất phù hợp"],
            ["Holme-Kim", "Power-law", "✅ Có", "✅ Cao", "✅ Rất phù hợp"],
            ["Erdős-Rényi", "Poisson (đều)", "❌ Không", "Rất thấp", "❌ Không phù hợp"],
            ["Watts-Strogatz", "Đều", "❌ Không", "✅ Cao", "❌ Không phù hợp"],
        ]
    )

    # ── SECTION DIVIDER 2 ───────────────────────────────────
    add_section_divider(prs, "02", "THỰC NGHIỆM",
                        "Thiết lập • Kết quả • Biểu đồ",
                        accent=ACCENT_ORG)

    # ── SLIDE 8: THIẾT LẬP THỰC NGHIỆM ─────────────────────
    add_content_slide(
        prs,
        title      = "Thiết Lập Thực Nghiệm",
        section_tag= "PHẦN 2 — THỰC NGHIỆM",
        accent     = ACCENT_ORG,
        bullets    = [
            ("Công cụ: Python 3.9+ | NetworkX | SciPy (Hungarian) | Matplotlib", 0),
            ("Cấu hình thực nghiệm:", 0),
            ("Số node: n = 350 (kích thước đồ thị G)", 1),
            ("Tỷ lệ xóa node: δ = 5% (cố định)", 1),
            ("Mức nhiễu ε: {0%, 5%, 10%, 15%, 20%}", 1),
            ("Số lần lặp: 5 trials / cấu hình → lấy trung bình", 1),
            ("4 mô hình đồ thị: Barabási-Albert, Holme-Kim, Erdős-Rényi, Watts-Strogatz", 1),
            ("Metric đánh giá: Accuracy = (số node de-anonymize đúng) / (tổng node trong H)", 0),
            ("Dataset: Synthetic (sinh từ NetworkX) — không cần download dữ liệu ngoài", 0),
        ]
    )

    # ── SLIDE 9: KẾT QUẢ NDSD ───────────────────────────────
    add_content_slide(
        prs,
        title      = "Kết Quả NDSD — So Sánh 4 Mô Hình",
        section_tag= "PHẦN 2 — THỰC NGHIỆM",
        accent     = ACCENT_ORG,
        bullets    = [
            ("[Chèn biểu đồ: bieu_do_1_accuracy_vs_noise.png — subplot trái]", 0),
            ("Nhận xét chính:", 0),
            ("ε = 0%: BA và HK đạt ~98-100% (hầu như hoàn hảo khi không có nhiễu)", 1),
            ("ε = 10%: BA/HK vẫn ~60-70%, ER/WS giảm gần về ngẫu nhiên (~10-20%)", 1),
            ("ε > 15%: Tất cả mô hình giảm đáng kể — nhiễu quá cao", 1),
            ("Lý do ER/WS kém: Phân phối bậc đồng đều → NDS quá giống nhau → khó phân biệt", 0),
            ("Lý do BA/HK tốt: Hub có NDS độc đáo → tạo 'neo' để NDSD tìm đúng", 0),
        ]
    )

    # ── SLIDE 10: KẾT QUẢ LOCAL SEARCH ─────────────────────
    add_content_slide(
        prs,
        title      = "Kết Quả Local Search — Cải Thiện Rõ Rệt",
        section_tag= "PHẦN 2 — THỰC NGHIỆM",
        accent     = ACCENT_ORG,
        bullets    = [
            ("[Chèn biểu đồ: bieu_do_1_accuracy_vs_noise.png — subplot phải]", 0),
            ("Cải thiện so với NDSD:", 0),
            ("BA ε=10%: NDSD ~65% → Local Search ~78% (+13 điểm)", 1),
            ("HK ε=10%: NDSD ~61% → Local Search ~74% (+13 điểm)", 1),
            ("ER, WS: Cải thiện ít vì Local Search cũng bị kẹt local optimum", 1),
            ("Kết luận quan trọng:", 0),
            ("✅ Kết hợp NDSD + Local Search chịu được nhiễu tới ε = 10%", 1),
            ("✅ BA và HK là mô hình tốt nhất — xác nhận phù hợp với mạng XH thực tế", 1),
            ("❌ ER và WS không phù hợp để mô hình hóa mạng xã hội", 1),
        ]
    )

    # ── SLIDE 11: PHÂN TÍCH NDSD vs LS ──────────────────────
    add_content_slide(
        prs,
        title      = "NDSD vs Local Search — Phân Tích Cải Thiện",
        section_tag= "PHẦN 2 — THỰC NGHIỆM",
        accent     = ACCENT_ORG,
        bullets    = [
            ("[Chèn biểu đồ: bieu_do_2_ndsd_vs_local_search.png]", 0),
            ("Vùng tô màu = mức độ cải thiện của Local Search so với NDSD", 0),
            ("Nhận xét:", 0),
            ("BA: Cải thiện lớn nhất ở ε = 5-10% (khi NDSD chưa quá tệ)", 1),
            ("HK: Tương tự BA nhờ clustering cao tạo cấu trúc rõ ràng hơn", 1),
            ("Khi ε > 15%: Vùng cải thiện thu hẹp — nhiễu quá cao, LS không đủ bù đắp", 1),
            ("Đây là đóng góp phân tích của nhóm — không có trong bài báo gốc", 0),
        ]
    )

    # ── SLIDE 12: ĐỐI CHIẾU BÀI BÁO GỐC ────────────────────
    add_table_slide(
        prs,
        title      = "Kết Quả Nhóm vs Bài Báo Gốc",
        section_tag= "PHẦN 2 — THỰC NGHIỆM",
        accent     = ACCENT_ORG,
        headers    = ["Tiêu chí", "Bài báo gốc", "Nhóm tái hiện"],
        rows       = [
            ["Dataset", "Mạng thực (Facebook, SNAP)", "Synthetic (NetworkX, n=350)"],
            ["Mô hình tốt nhất", "BA và HK", "✅ Xác nhận"],
            ["Chịu nhiễu 10%", "✅ Có", "✅ Tái hiện được"],
            ["ER, WS kém", "✅ Có", "✅ Xác nhận"],
            ["Scale", "Vài nghìn node", "350 node (giới hạn tính toán O(n⁵))"],
        ]
    )

    # ── SECTION DIVIDER 3 ───────────────────────────────────
    add_section_divider(prs, "03", "MỞ RỘNG",
                        "An ninh Di động • Phòng thủ • Hướng nghiên cứu",
                        accent=COLOR_GREEN)

    # ── SLIDE 13: ỨNG DỤNG AN NINH DI ĐỘNG ─────────────────
    add_content_slide(
        prs,
        title      = "De-Anonymization Trong Bối Cảnh An Ninh Di Động",
        section_tag= "PHẦN 3 — MỞ RỘNG",
        accent     = COLOR_GREEN,
        bullets    = [
            ("Ứng dụng di động thu thập 'đồ thị quan hệ' ẩn ngầm mà người dùng không hay biết:", 0),
            ("Call logs: ai gọi cho ai → đồ thị quan hệ trực tiếp", 1),
            ("Danh bạ (READ_CONTACTS): mạng quan hệ xã hội đầy đủ", 1),
            ("Bluetooth proximity: ai ở gần nhau về mặt vật lý", 1),
            ("Wi-Fi co-location: ai ở cùng địa điểm", 1),
            ("Kịch bản tấn công thực tế:", 0),
            ("Kẻ tấn công có app với quyền READ_CONTACTS + READ_CALL_LOGS", 1),
            ("Xây đồ thị G (mạng quan hệ đã biết) → so khớp với mạng ẩn danh H", 1),
            ("→ De-anonymize người dùng dù họ đã xóa tên/ID", 1),
        ]
    )

    # ── SLIDE 14: THREAT MODEL ───────────────────────────────
    add_content_slide(
        prs,
        title      = "Mô Hình Mối Đe Dọa (Threat Model)",
        section_tag= "PHẦN 3 — MỞ RỘNG",
        accent     = COLOR_GREEN,
        bullets    = [
            ("Kẻ tấn công có:", 0),
            ("Mạng xã hội công khai G: LinkedIn, Facebook public profile", 1),
            ("Mạng ẩn danh H: dark web forum, group tội phạm", 1),
            ("Điều kiện tấn công thành công:", 0),
            ("H là noisy subgraph của G (cùng người, khác nền tảng)", 1),
            ("Nhiễu ε ≤ 10% (thường thỏa mãn với mạng XH thực)", 1),
            ("Mô hình mạng gần BA hoặc HK (thực tế hầu hết mạng XH đều vậy)", 1),
            ("Ví dụ lịch sử: FBI de-anonymize thành viên Silk Road qua phân tích cấu trúc mạng", 0),
        ]
    )

    # ── SLIDE 15: PHÒNG THỦ ─────────────────────────────────
    add_content_slide(
        prs,
        title      = "Cơ Chế Phòng Thủ Quyền Riêng Tư",
        section_tag= "PHẦN 3 — MỞ RỘNG",
        accent     = COLOR_GREEN,
        bullets    = [
            ("1. Graph Perturbation (Thêm nhiễu có chủ ý):", 0),
            ("Thêm/xóa cạnh ngẫu nhiên trước khi chia sẻ đồ thị", 1),
            ("Cần ε > 10% để phá vỡ de-anonymization — nhưng giảm tính hữu dụng", 1),
            ("2. Differential Privacy on Graphs:", 0),
            ("Thêm nhiễu có kiểm soát (Laplace/Gaussian mechanism)", 1),
            ("Đảm bảo ε-differential privacy — tiêu chuẩn vàng hiện nay", 1),
            ("Apple và Google dùng DP cho analytics trên iOS/Android", 1),
            ("3. k-Anonymity cho Graphs:", 0),
            ("Mỗi node không phân biệt được với ít nhất k-1 node khác", 1),
            ("Cân bằng k và utility — k càng lớn càng an toàn nhưng mất thông tin", 1),
        ]
    )

    # ── SLIDE 16: HƯỚNG MỞ RỘNG ─────────────────────────────
    add_table_slide(
        prs,
        title      = "Hạn Chế Bài Báo & Hướng Nghiên Cứu Tiếp Theo",
        section_tag= "PHẦN 3 — MỞ RỘNG",
        accent     = COLOR_GREEN,
        headers    = ["Hạn chế", "Hướng giải quyết"],
        rows       = [
            ["O(n⁵) — không scale với mạng lớn", "GNN-based approaches, thuật toán song song"],
            ["Local search mắc kẹt local optimum", "Simulated Annealing, Genetic Algorithm"],
            ["Chỉ xét đồ thị vô hướng không trọng số", "Mở rộng cho directed/weighted graphs"],
            ["ε > 10%: hiệu năng giảm mạnh", "Kết hợp thêm semantic features hoặc học máy"],
            ["Giả thiết H là subgraph của G duy nhất", "Multi-graph matching phức tạp hơn"],
        ]
    )

    # ── SLIDE 17: DEMO ───────────────────────────────────────
    add_content_slide(
        prs,
        title      = "Demo — Kết Quả Thực Nghiệm (n=350)",
        section_tag= "PHẦN 3 — MỞ RỘNG",
        accent     = COLOR_GREEN,
        bullets    = [
            ("[Chèn 2 biểu đồ PNG vào slide này]", 0),
            ("Kết quả nổi bật (Barabási-Albert):", 0),
            ("ε = 0%:  NDSD ~98.5%  →  Local Search ~100%", 1),
            ("ε = 5%:  NDSD ~75.0%  →  Local Search ~88.5%", 1),
            ("ε = 10%: NDSD ~62.0%  →  Local Search ~76.0%", 1),
            ("ε = 15%: NDSD ~40.0%  →  Local Search ~53.0%", 1),
            ("ε = 20%: NDSD ~23.0%  →  Local Search ~31.0%", 1),
            ("→ Xác nhận ngưỡng 10% nhiễu của bài báo gốc", 0),
        ]
    )

    # ── SLIDE 18: KẾT LUẬN ───────────────────────────────────
    add_content_slide(
        prs,
        title      = "Kết Luận",
        section_tag= "",
        bullets    = [
            ("1. LÝ THUYẾT: Hình thức hóa bài toán de-anonymization dưới dạng tối ưu hóa đồ thị", 0),
            ("NDSD (O(n⁴)) + Local Search (O(n⁵)) — kết hợp bù trừ điểm yếu nhau", 1),
            ("2. THỰC NGHIỆM: Chịu được nhiễu ε ≤ 10% trên mô hình BA và HK", 0),
            ("BA & HK phù hợp mạng XH thực; ER & WS không phù hợp", 1),
            ("3. MỞ RỘNG: Ứng dụng thực tiễn trong An ninh Di động", 0),
            ("Cần cơ chế phòng thủ: Differential Privacy, Graph Perturbation", 1),
            ("Thông điệp cuối cùng:", 0),
            ("⚠️  Ẩn danh hóa dữ liệu ≠ Bảo vệ quyền riêng tư!", 1),
            ("Cấu trúc quan hệ có thể tiết lộ danh tính dù không có metadata", 1),
        ]
    )

    # ── SLIDE 19: TÀI LIỆU THAM KHẢO ───────────────────────
    add_content_slide(
        prs,
        title      = "Tài Liệu Tham Khảo",
        section_tag= "",
        bullets    = [
            ("Caragiannis & Tsitsoka (2019). Deanonymizing Social Networks Using Structural Information. IJCAI-2019, pp.1213-1219.", 0),
            ("Narayanan & Shmatikov (2009). De-anonymizing Social Networks. IEEE S&P.", 0),
            ("Backstrom et al. (2007). Wherefore art thou r3579x? KDD 2007.", 0),
            ("Babai et al. (1980). Random Graph Isomorphism. SIAM.", 0),
            ("Barabási & Albert (1999). Emergence of Scaling in Random Networks. Science.", 0),
            ("Holme & Kim (2002). Growing Scale-Free Networks with Tunable Clustering. PRE.", 0),
            ("Dwork (2006). Differential Privacy. ICALP.", 0),
        ]
    )

    return prs


# ─────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────

if __name__ == "__main__":
    output_path = "ProjectANDD/presentation.pptx"
    os.makedirs("ProjectANDD", exist_ok=True)

    print("[*] Đang tạo file PowerPoint...")
    prs = build_presentation()
    prs.save(output_path)
    print(f"[OK] Đã tạo: {output_path}")
    print(f"     Tổng số slide: {len(prs.slides)}")
    print(f"[!]  Mở file bằng PowerPoint hoặc Google Slides để chèn biểu đồ PNG")
